import os
from wand.image import Image 
from pptx import Presentation  
from pptx.util import Inches 


# Creating a presentation object
ppt = Presentation()

# Resizing the presentation slide witdth
ppt.slide_width = Inches(13)

# Position for image in each slide
left =Inches(1) 
top = Inches(2.5)

# Height of the image
height = Inches(4.5)

# Resizing the Logo image

logo = Image(filename=f"image1/nike_black.png") # image filefolder for inserting logo
logo_width = 1200
logo_height = 400
logo.resize(logo_width, logo_height)

# Position coordinates for logo image
x = 100 
y = 80

# Initialize the for loop
for i in range(1,6):
    with Image(filename = f'image1/image{i}.jpg' ) as image:
        
        image.composite(logo,x , y) # Appending the logo to the image
        image_name = f'output{i}.jpg' # Name of the output image for future reference
        image.save(filename = image_name ) # saveing final output image


        # Adding Slide to ppt
        slide = ppt.slides.add_slide(ppt.slide_layouts[1]) 
        slide.shapes.title.text = " ImageMagick python-ppt" # Title for the slide
        slide.placeholders[1].text = image_name # Subtitle for the slide

        # Adding Image to slide
        pic = slide.shapes.add_picture(image_name,left, top,height = height)


# Saving the final result 
ppt.save('output.pptx')